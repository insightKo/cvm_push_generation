# -*- coding: utf-8 -*-
"""Роадмэп и статус проекта внедрения ИИ в ДИКСИ на 02.09.2026 (PPTX).
Источник плана и трудоёмкости: data/Оценка_трудоёмкости_внедрение_ДИКСИ_06.08.2026 (1).xlsx
(лист «Трудоёмкость (детально)» — помесячная раскладка задач; лист «Roadmap» — эффект).
Стиль: bonnie-slide + mckinsey-deck, тело >=14pt, 10-11pt только в плотной диаграмме Ганта."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

INK    = RGBColor(0x1E, 0x1E, 0x2F)
ORANGE = RGBColor(0xEE, 0x72, 0x03)
PURPLE = RGBColor(0x6C, 0x4A, 0xB6)
GREEN  = RGBColor(0x2E, 0x9E, 0x5B)
GREY   = RGBColor(0x5B, 0x61, 0x6E)
LGREY  = RGBColor(0xF4, 0xF5, 0xF7)
BORD   = RGBColor(0xD9, 0xDC, 0xE1)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
WARM   = RGBColor(0xFD, 0xF8, 0xF2)
MUTE   = RGBColor(0xC9, 0xCD, 0xD6)
PLAN   = RGBColor(0xB9, 0xBF, 0xCA)
GREENBG = RGBColor(0xEF, 0xF6, 0xF0)
ORNGBG = RGBColor(0xFD, 0xF1, 0xE4)
FONT   = "Arial"
SW, SH = 13.333, 7.5

# ── реальный перенос строк по метрикам Arial (а не по «примерно N символов»)
from PIL import ImageFont as _IF
_ARIAL = "/System/Library/Fonts/Supplemental/Arial.ttf"
_fcache = {}


def _font(pt, bold=False):
    key = (round(pt * 4), bold)
    if key not in _fcache:
        path = _ARIAL.replace("Arial.ttf", "Arial Bold.ttf") if bold else _ARIAL
        _fcache[key] = _IF.truetype(path, int(round(pt * 4)))  # 4x для точности
    return _fcache[key]


def text_w(t, pt, bold=False):
    """Ширина строки в дюймах при кегле pt."""
    return _font(pt, bold).getlength(t) / 4.0 / 72.0


def wrap_lines(t, w_in, pt, bold=False):
    """Число строк при переносе по словам в колонке шириной w_in дюймов."""
    words, lines, cur = t.split(), 0, ""
    for w in words:
        trial = (cur + " " + w).strip()
        if text_w(trial, pt, bold) <= w_in or not cur:
            cur = trial
        else:
            lines += 1; cur = w
    return lines + (1 if cur else 0)


DONE, WIP, PLND = "done", "wip", "plan"
ST_COL = {DONE: GREEN, WIP: ORANGE, PLND: PLAN}
ST_TXT = {DONE: "Сделано", WIP: "В работе", PLND: "План"}


def new_prs():
    p = Presentation(); p.slide_width = Inches(SW); p.slide_height = Inches(SH)
    return p


def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def box(s, l, t, w, h, fill=WHITE, line=BORD, lw=0.75, rounded=True):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    if rounded:
        try: shp.adjustments[0] = 0.045
        except Exception: pass
    return shp


def txt(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space_after=2, ls=1.0):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for m in ("left", "right"): setattr(tf, f"margin_{m}", Inches(0.05))
    for m in ("top", "bottom"): setattr(tf, f"margin_{m}", Inches(0.02))
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(0); p.line_spacing = ls
        for (t_, sz, col, bold) in para:
            r = p.add_run(); r.text = t_
            r.font.size = Pt(sz); r.font.color.rgb = col; r.font.bold = bold; r.font.name = FONT
    return tb


def header_block(s, kicker, title, page):
    box(s, 0.5, 0.4, 0.1, 0.82, fill=ORANGE, line=None, rounded=False)
    txt(s, 0.74, 0.38, 12.1, 0.32, [[(kicker.upper(), 12.5, ORANGE, True)]])
    txt(s, 0.74, 0.66, 12.2, 0.72, [[(title, 21, INK, True)]])
    txt(s, 0.5, 7.04, 9.5, 0.32,
        [[("Внедрение ИИ в CVM · ДИКСИ · статус на 02.09.2026", 10, GREY, False)]])
    txt(s, 12.0, 7.04, 0.85, 0.32, [[(str(page), 10, GREY, True)]], align=PP_ALIGN.RIGHT)


def card(s, l, t, w, h, lines, fill=WHITE, line=BORD, accent=None, ls=1.0):
    box(s, l, t, w, h, fill=fill, line=line, lw=0.75)
    if accent:
        box(s, l, t, 0.07, h, fill=accent, line=None, rounded=False)
    txt(s, l + 0.18, t + 0.03, w - 0.26, h - 0.06, lines, anchor=MSO_ANCHOR.MIDDLE, ls=ls, space_after=2)


def legend(s, l, t):
    x = l
    for st in (DONE, WIP, PLND):
        box(s, x, t + 0.045, 0.26, 0.16, fill=ST_COL[st], line=None, rounded=False)
        txt(s, x + 0.34, t, 1.15, 0.26, [[(ST_TXT[st], 11.5, GREY, False)]])
        x += 1.42
    return x


# ── месячная ось: авг.26 … мар.27 (как в листе «Трудоёмкость (детально)»)
MONTHS = ["авг.26", "сент.26", "окт.26", "ноя.26", "дек.26", "янв.27", "фев.27", "мар.27"]


def gantt(s, rows, gx, gy, gw, row_h, name_w, today_idx=1.9, milestone=None, name_sz=11.5):
    """rows: (label, sub, m_from, m_to, status) — m_from/m_to индексы месяцев (вкл.)."""
    cw = (gw - name_w) / len(MONTHS)
    # шапка месяцев
    for i, m in enumerate(MONTHS):
        x = gx + name_w + i * cw
        box(s, x, gy, cw, 0.34, fill=INK if i < 3 else GREY, line=WHITE, lw=1.0, rounded=False)
        txt(s, x, gy, cw, 0.34, [[(m, 11, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    n = len(rows)
    ty = gy + 0.34
    # фон строк + вертикальные линии месяцев
    for j in range(n):
        box(s, gx, ty + j * row_h, gw, row_h, fill=WHITE if j % 2 else LGREY, line=None, rounded=False)
    for i in range(len(MONTHS) + 1):
        x = gx + name_w + i * cw
        ln = box(s, x, ty, 0.008, n * row_h, fill=BORD, line=None, rounded=False)
    # линия «сегодня»
    tx = gx + name_w + today_idx * cw
    box(s, tx, gy, 0.022, 0.34 + n * row_h, fill=ORANGE, line=None, rounded=False)
    txt(s, tx - 0.72, gy - 0.28, 1.44, 0.26, [[("02.09 · 36 неделя", 10, ORANGE, True)]], align=PP_ALIGN.CENTER)
    if milestone:
        mi, mlabel = milestone
        mx = gx + name_w + mi * cw
        d = s.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(mx - 0.09), Inches(ty + n * row_h + 0.05),
                               Inches(0.18), Inches(0.18))
        d.fill.solid(); d.fill.fore_color.rgb = PURPLE; d.line.fill.background(); d.shadow.inherit = False
        txt(s, mx + 0.14, ty + n * row_h + 0.02, 5.3, 0.26, [[(mlabel, 11, PURPLE, True)]])
    # строки
    for j, (label, sub, a, b, st) in enumerate(rows):
        y = ty + j * row_h
        if label.startswith("§"):   # заголовок этапа
            box(s, gx, y, gw, row_h, fill=WARM, line=None, rounded=False)
            txt(s, gx + 0.1, y, gw - 0.2, row_h, [[(label[1:], 12, ORANGE, True)]], anchor=MSO_ANCHOR.MIDDLE)
            continue
        runs = [[(label, name_sz, INK, False)]]
        if sub:
            runs[0].append(("  " + sub, name_sz - 1.5, GREY, False))
        txt(s, gx + 0.1, y, name_w - 0.16, row_h, runs, anchor=MSO_ANCHOR.MIDDLE, ls=0.95)
        bx = gx + name_w + a * cw + 0.05
        bw = (b - a + 1) * cw - 0.10
        bar = box(s, bx, y + row_h * 0.22, bw, row_h * 0.56, fill=ST_COL[st], line=None, rounded=True)
        txt(s, bx, y + row_h * 0.22, bw, row_h * 0.56, [[(ST_TXT[st], 9.5, WHITE, True)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    frame = box(s, gx, gy, gw, 0.34 + n * row_h, fill=WHITE, line=BORD, rounded=False)
    frame.fill.background()


# ════════════════════════════ Слайд 1 — общий статус
def build_overview(prs):
    s = add_slide(prs)
    header_block(s, "Роадмэп · статус на 02.09.2026",
                 "Этап 1 закрыт на 6 задач из 13: акции ДЦО на согласовании в КД, запуск коммуникаций — 37 неделя (07.09.–12.09.)", 1)
    # KPI-плашки
    tiles = [
        ("ЭТАП 1 · ДЦО", "6 из 13", "задач этапа закрыто", GREEN),
        ("СЕНТЯБРЬ · ДЦО", "+43 млн ₽", "доп. ТО, расчёт от 02.09\n(в плане было 27 млн ₽)", ORANGE),
        ("СТАТУС АКЦИЙ", "На согласовании", "акции ДЦО высланы в КД", PURPLE),
        ("ПЕРВЫЕ КОММУНИКАЦИИ", "37 неделя", "07.09.–12.09., старт пилота\nс контрольными группами", INK),
    ]
    tw, gap, tl, tt, th = 3.02, 0.22, 0.6, 1.42, 1.30
    for i, (k, big, sub, col) in enumerate(tiles):
        x = tl + i * (tw + gap)
        box(s, x, tt, tw, th, fill=WHITE, line=BORD)
        box(s, x, tt, tw, 0.09, fill=col, line=None, rounded=False)
        txt(s, x + 0.18, tt + 0.16, tw - 0.3, 0.26, [[(k, 10.5, col, True)]])
        txt(s, x + 0.18, tt + 0.44, tw - 0.3, 0.40, [[(big, 21, INK, True)]])
        txt(s, x + 0.18, tt + 0.85, tw - 0.3, 0.42,
            [[(l, 12, GREY, False)] for l in sub.split("\n")], ls=1.0)

    # Полоса этапов
    txt(s, 0.6, 3.00, 12.2, 0.3, [[("ЧЕТЫРЕ ЭТАПА · АВГУСТ 2026 — МАРТ 2027", 12, ORANGE, True)]])
    stages = [
        ("ЭТАП 1 · авг–окт 26", "Персонализация ДЦО", "670 ч", DONE,
         "Сервис, данные, MCI, КГ,\nдвижок текстов, контент ДЦО"),
        ("ЭТАП 2 · ноя 26 – янв 27", "CVM: 100 сегментов", "1 202 ч", PLND,
         "100 сегментов, конвейер контента,\nинтерфейс, агенты, дашборды"),
        ("ЭТАП 3 · фев 27", "Масштаб 200 акций", "310 ч", PLND,
         "Оптимизатор, 200 акций,\nпередача сервиса"),
        ("ЭТАП 4 · мар 27", "Дообучение на откликах", "299 ч", PLND,
         "Витрина откликов, дообучение,\nстоп-листы, итоги пилота"),
    ]
    sw_, sg, sl, st_, sh_ = 2.98, 0.24, 0.6, 3.34, 1.60
    for i, (per, name, hrs, status, det) in enumerate(stages):
        x = sl + i * (sw_ + sg)
        fill = GREENBG if status == DONE else WHITE
        box(s, x, st_, sw_, sh_, fill=fill, line=BORD)
        box(s, x, st_, sw_, 0.07, fill=ST_COL[status], line=None, rounded=False)
        txt(s, x + 0.18, st_ + 0.14, sw_ - 0.32, 0.26, [[(per, 10.5, GREY, True)]])
        txt(s, x + 0.18, st_ + 0.40, sw_ - 0.32, 0.34, [[(name, 15, INK, True)]])
        txt(s, x + 0.18, st_ + 0.78, sw_ - 0.32, 0.60,
            [[(l, 11, GREY, False)] for l in det.split("\n")], ls=1.05)
        txt(s, x + 0.18, st_ + sh_ - 0.34, sw_ - 0.32, 0.26,
            [[(hrs + " · ", 11.5, GREY, True), (ST_TXT[status] if status != DONE else "Этап в работе", 11.5, ST_COL[status], True)]])
        if i < 3:
            a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + sw_ + 0.045), Inches(st_ + 0.66),
                                   Inches(0.15), Inches(0.30))
            a.fill.solid(); a.fill.fore_color.rgb = ORANGE; a.line.fill.background(); a.shadow.inherit = False

    # Нижняя лента — эффект (ДЦО + CVM, стопка)
    box(s, 0.6, 5.20, 12.14, 1.66, fill=WARM, line=BORD)
    txt(s, 0.8, 5.28, 4.2, 0.3, [[("ДОП. ТОВАРООБОРОТ, млн ₽/мес", 12, ORANGE, True)]])
    # легенда
    lx = 5.25
    for lbl, col in (("ДЦО", ORANGE), ("CVM", PURPLE)):
        box(s, lx, 5.34, 0.22, 0.14, fill=col, line=None, rounded=False)
        txt(s, lx + 0.28, 5.26, 0.75, 0.26, [[(lbl, 11.5, GREY, False)]])
        lx += 1.02
    # сент/окт: ячейки CVM в файле пустые — показываем только ДЦО
    eff = [("сент.26", 43, None), ("окт.26", 54, None), ("ноя.26", 54, 208),
           ("дек.26", 34, 208), ("янв.27", 54, 100), ("фев.27", 54, 150)]
    bx0, byb, bwd, maxh, mx = 0.85, 6.52, 0.80, 0.70, 262.0
    for i, (m, dts, cvm) in enumerate(eff):
        x = bx0 + i * 1.02
        hd = maxh * dts / mx
        box(s, x, byb - hd, bwd, hd, fill=ORANGE, line=None, rounded=False)
        top = byb - hd
        if cvm:
            # в стопке подписываем каждую часть внутри своего сегмента
            txt(s, x, byb - hd, bwd, hd, [[(str(dts), 10, WHITE, True)]],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            hc = maxh * cvm / mx
            box(s, x, top - hc, bwd, hc, fill=PURPLE, line=None, rounded=False)
            txt(s, x, top - hc, bwd, hc, [[(str(cvm), 11, WHITE, True)]],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            top -= hc
        total = dts + (cvm or 0)
        txt(s, x - 0.14, top - 0.25, bwd + 0.28, 0.24,
            [[(str(total), 12, INK, True)]], align=PP_ALIGN.CENTER)
        txt(s, x - 0.14, byb + 0.03, bwd + 0.28, 0.24, [[(m, 10.5, GREY, False)]], align=PP_ALIGN.CENTER)
    txt(s, 7.30, 5.62, 5.3, 1.05,
        [[("ДЦО за период — 243 млн ₽,  CVM — 1 040 млн ₽", 13.5, INK, True)],
         [("Совокупно за 6 месяцев ~1,28 млрд ₽", 13.5, INK, True)],
         [("Сентябрь по ДЦО пересчитан: 43 млн ₽ вместо 27 млн ₽ в плане", 12, ORANGE, False)],
         [("CVM включается с этапа 2. В файле оценки за сент.–окт. значение CVM не проставлено.", 10.5, GREY, False)]], ls=1.1)


# ════════════════════════════ Слайд 2 — что сделано на 02.09
def build_status(prs):
    s = add_slide(prs)
    header_block(s, "Статус · 2 сентября 2026",
                 "Сделано на 2 сентября: сервис, данные ДЦО, сегментация MCI, контрольные группы, контент акций", 2)
    col_w, gap, ty = 3.94, 0.22, 1.44
    ch = 5.28
    cols = [
        ("СДЕЛАНО", GREEN, GREENBG, [
            "Сервис развёрнут: доступы, безопасность, бэкапы",
            "Подключены данные ДЦО: клиенты, покупки, текущие акции",
            "Встроен модуль сегментации MCI + интерфейс выбора сегментов",
            "Настроены автоматические контрольные группы на базе MCI",
            "Описана методология «сегмент → оффер»",
            "ДЦО-акции переведены в персональный формат",
        ]),
        ("В РАБОТЕ", ORANGE, ORNGBG, [
            "Акции ДЦО высланы на согласование в КД",
            "Движок персонализации текстов ДЦО под сегмент",
            "Агенты в обучении: копирайтер, редактор-гуманизатор",
            "Первые коммуникации — 37 неделя (07.09.–12.09.)",
            "База знаний: шаблоны, скиллы, гайдлайны, tone of voice",
            "Бизнес-процесс планирования и заведения ДЦО",
            "Отдельные среды под агентов, шифрование данных",
            "Пересчитан эффект сентября: +43 млн ₽ по ДЦО",
        ]),
        ("СЛЕДУЮЩИЕ ШАГИ", PURPLE, LGREY, [
            "Запуск коммуникаций и замер vs контрольные группы",
            "Постэффекты пилота: отклик, доп. ТО, маржа по сегментам",
            "Приёмка этапа 1, старт этапа 2 — 100 сегментов с ноября",
            "Конвейер контента и интерфейс заведения акций",
            "Агенты: конструктор акций, оператор заведения, дизайнер",
        ]),
    ]
    for i, (title_, col, bg, items) in enumerate(cols):
        x = 0.6 + i * (col_w + gap)
        box(s, x, ty, col_w, ch, fill=bg, line=BORD)
        box(s, x, ty, col_w, 0.42, fill=col, line=None, rounded=False)
        txt(s, x + 0.18, ty, col_w - 0.3, 0.42, [[(title_, 13, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
        iy = ty + 0.58
        tw = col_w - 0.62
        for it in items:
            n_lines = wrap_lines(it, tw - 0.04, 14)
            h = 0.215 * n_lines
            box(s, x + 0.20, iy + 0.075, 0.12, 0.12, fill=col, line=None, rounded=False)
            txt(s, x + 0.44, iy - 0.03, tw, h + 0.1, [[(it, 14, INK, False)]], ls=1.02)
            iy += h + 0.14


# ════════════════════════════ Слайд 3 — Гант этапа 1
def build_gantt_stage1(prs):
    s = add_slide(prs)
    header_block(s, "Роадмэп · этап 1",
                 "Персонализация ДЦО: данные, сегментация и контент готовы; движок персонализации, пилот и процесс — в работе", 3)
    rows = [
        ("Развёртывание сервиса: доступы, безопасность, бэкапы", "1.1", 0, 1, DONE),
        ("Интеграция источников данных ДЦО", "1.2", 0, 1, DONE),
        ("Модуль сегментации MCI + интерфейс сегментов", "1.3", 0, 1, DONE),
        ("Автоматические контрольные группы на базе MCI", "1.4", 0, 1, DONE),
        ("Движок персонализации текстов ДЦО", "1.5", 0, 2, WIP),
        ("Методология персонализации: сегмент → оффер", "1.6", 0, 1, DONE),
        ("Перевод ДЦО-акций в персональный формат", "1.7", 0, 1, DONE),
        ("База знаний: шаблоны, скиллы, tone of voice", "1.8", 1, 2, WIP),
        ("Бизнес-процесс планирования и заведения ДЦО", "1.9", 1, 2, WIP),
        ("Пилот + постэффекты vs контрольная группа", "1.10", 1, 2, WIP),
        ("Серверы под агентов: изолированные среды", "1.12", 1, 2, WIP),
        ("Безопасность: шифрование данных для агентов", "1.13", 1, 2, WIP),
        ("Управление этапом: статусы, демо, приёмка", "1.11", 0, 2, WIP),
    ]
    gantt(s, rows, gx=0.6, gy=1.70, gw=12.14, row_h=0.33, name_w=5.55,
          milestone=(1.55, "Первые персональные коммуникации — 37 неделя (07.09.–12.09.)"))
    legend(s, 0.62, 1.34)
    txt(s, 5.6, 1.30, 7.2, 0.3,
        [[("Этап 1 · 13 задач: 6 закрыто, 7 в работе", 12.5, GREY, False)]],
        align=PP_ALIGN.RIGHT)
    txt(s, 0.6, 6.72, 12.2, 0.3,
        [[("Акции ДЦО по задаче 1.7 переданы на согласование в коммерческую дирекцию; пилот 1.10 стартует после согласования.", 12.5, GREY, False)]])


# ════════════════════════════ Слайд 4 — Гант этапов 2–4
def build_gantt_next(prs):
    s = add_slide(prs)
    header_block(s, "Роадмэп · этапы 2–4",
                 "С ноября — масштаб до 100 сегментов, затем 200 акций в месяц и дообучение агентов на откликах", 4)
    rows = [
        ("§ЭТАП 2 · ноябрь 2026 — январь 2027 · Полноценный CVM: 100 сегментов · 1 202 ч", "", 0, 0, PLND),
        ("100 сегментов MCI в проекте: конфиг, обновление, выгрузки", "2.1", 3, 5, PLND),
        ("Конвейер генерации контента + контент под 100 сегментов", "2.2 · 2.7", 3, 5, PLND),
        ("Операционный интерфейс заведения акций", "2.3", 3, 5, PLND),
        ("Деплинки, SKU и механики в конвейере", "2.4", 3, 5, PLND),
        ("Агенты: конструктор, оператор, дизайнер, оптимизатор", "2.5 · 2.13 · 2.15–2.16", 3, 5, PLND),
        ("Дашборд эффективности и постэффекты сегментных акций", "2.6 · 2.14", 3, 5, PLND),
        ("Методология 100 сегментов, процесс, обучение команды заказчика", "2.8–2.12", 3, 5, PLND),
        ("§ЭТАП 3 · февраль 2027 · Масштаб + простая оптимизация · 310 ч", "", 0, 0, PLND),
        ("200 акций через конвейер + автоподбор оффера и сегмента", "3.2 · 3.3", 6, 6, PLND),
        ("Стабилизация интерфейса, отчёт по эффекту, передача сервиса", "3.1 · 3.4–3.6", 6, 6, PLND),
        ("§ЭТАП 4 · март 2027 · Дообучение на откликах и ограничения товаров · 299 ч", "", 0, 0, PLND),
        ("Витрина откликов клиентов, дообучение и калибровка агентов", "4.1–4.4", 7, 7, PLND),
        ("Стоп-листы товаров и ограничения в конструкторе акций", "4.5 · 4.6", 7, 7, PLND),
        ("Коррекция методологии по итогам пилота, итоговая аналитика", "4.7–4.10", 7, 7, PLND),
    ]
    gantt(s, rows, gx=0.6, gy=1.70, gw=12.14, row_h=0.30, name_w=6.05, name_sz=11)
    legend(s, 0.62, 1.34)


def main():
    prs = new_prs()
    build_overview(prs)
    build_status(prs)
    build_gantt_stage1(prs)
    build_gantt_next(prs)
    out = "output/Роадмэп_статус_внедрение_ИИ_ДИКСИ_02.09.2026.pptx"
    prs.save(out)
    print("saved:", out)


if __name__ == "__main__":
    main()
